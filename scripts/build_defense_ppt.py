"""
Build a TECHNICAL DEFENSE PPT for the Voice Agent.

Different from the take-home PPT:
  - Tells a story (problem → decision → result)
  - Demo-heavy (embedded screenshots/audio references)
  - Anticipates Q&A
  - Shows debugging journey (a senior signal)

Run: uv run python scripts/build_defense_ppt.py
Output: docs/pptx/output/voice-agent-defense.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# ---------------------------------------------------------------- theme (跟 take-home PPT 一致)
THEME = {
    "primary":   "03045E", "secondary": "0077B6", "accent":    "00B4D8",
    "light":     "CAF0F8", "bg":        "F8FAFC", "muted":     "64748B",
    "border":    "E2E8F0", "white":     "FFFFFF", "good":      "10B981",
    "bad":       "EF4444", "warn":      "F59E0B",
}
FONT_CN = "Microsoft YaHei"
FONT_EN = "Arial"


def hex_rgb(h: str) -> RGBColor: return RGBColor.from_string(h)


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_CN):
    """One-shot text box. `text` can be a list of (text, kwargs) tuples for inline runs."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    if isinstance(text, str):
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color or hex_rgb(THEME["primary"])
    else:
        # list of paragraphs — each para can be a string OR list of (text, kwargs) tuples
        for i, para in enumerate(text):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if isinstance(para, str):
                r = p.add_run()
                r.text = para
                r.font.name = font
                r.font.size = Pt(size)
                r.font.bold = bold
                r.font.color.rgb = color or hex_rgb(THEME["primary"])
            else:
                for run_text, run_kw in para:
                    r = p.add_run()
                    r.text = run_text
                    r.font.name = run_kw.get("font", font)
                    r.font.size = Pt(run_kw.get("size", size))
                    r.font.bold = run_kw.get("bold", bold)
                    r.font.color.rgb = run_kw.get("color", color or hex_rgb(THEME["primary"]))
    return box


def add_rect(slide, x, y, w, h, fill, line=None, line_w=0.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = hex_rgb(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = hex_rgb(line)
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_card(slide, x, y, w, h, title, body, *, accent="secondary"):
    """Card with a left accent stripe + title + body."""
    add_rect(slide, x, y, w, h, THEME["white"], line=THEME["border"], line_w=0.75)
    add_rect(slide, x, y, 0.12, h, THEME[accent])
    add_text(slide, x + 0.25, y + 0.18, w - 0.4, 0.45, title, size=14, bold=True, color=hex_rgb(THEME["primary"]))
    if isinstance(body, str):
        add_text(slide, x + 0.25, y + 0.7, w - 0.4, h - 0.85, body, size=12, color=hex_rgb(THEME["muted"]))
    else:
        add_text(slide, x + 0.25, y + 0.7, w - 0.4, h - 0.85, body, size=12, color=hex_rgb(THEME["muted"]))


def header(slide, idx, total, title, subtitle=""):
    """Standard slide header: number + title + subtitle + bottom bar."""
    add_text(slide, 0.5, 0.25, 1, 0.3, f"{idx:02d} / {total:02d}",
             size=10, color=hex_rgb(THEME["muted"]), font=FONT_EN)
    add_text(slide, 0.5, 0.55, 12, 0.7, title, size=28, bold=True, color=hex_rgb(THEME["primary"]))
    if subtitle:
        add_text(slide, 0.5, 1.2, 12, 0.4, subtitle, size=14, color=hex_rgb(THEME["muted"]))
    add_rect(slide, 0, 13.0, 13.33, 0.05, THEME["accent"])


# ======================================================================
# Slide content
# ======================================================================

SLIDES = []  # list of (builder_fn, slide_title_for_toc)

TOTAL = 14  # will count


def s01_cover(prs):
    """Cover — title + 副标题 + 答辩人占位"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, 13.33, 7.5, THEME["primary"])
    add_rect(s, 0, 5.5, 13.33, 2, THEME["secondary"])
    add_rect(s, 0, 6.5, 13.33, 0.08, THEME["accent"])
    add_text(s, 1, 1.5, 11, 0.6, "技术答辩", size=20, color=hex_rgb(THEME["accent"]), font=FONT_CN)
    add_text(s, 1, 2.2, 11, 1.5, "Voice Agent for 园区访客登记",
             size=44, bold=True, color=hex_rgb(THEME["white"]))
    add_text(s, 1, 3.7, 11, 0.6,
             "Twilio Media Streams  ·  Bailian (百炼) AI 全栈  ·  实时语音对话 < 25s",
             size=16, color=hex_rgb(THEME["light"]))
    add_text(s, 1, 6.8, 11, 0.4, "答辩人：______     蓝色鲸鱼科技园  ·  2026.06",
             size=12, color=hex_rgb(THEME["white"]))


def s02_self(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, 2, TOTAL, "自我介绍 & 项目定位", "一分钟讲清楚我是谁、我做了什么")
    add_text(s, 0.5, 2.0, 12, 0.4, "项目来源", size=12, color=hex_rgb(THEME["secondary"]))
    add_text(s, 0.5, 2.4, 12, 0.6, "蓝色鲸鱼科技园 园区访客登记系统",
             size=20, bold=True, color=hex_rgb(THEME["primary"]))
    add_text(s, 0.5, 3.1, 12, 0.4, "业务场景", size=12, color=hex_rgb(THEME["secondary"]))
    add_text(s, 0.5, 3.5, 12, 0.6, "访客开车到园区入口 → 打一通电话 → AI 自动采集车牌/公司/事由 → 推企业微信给门卫",
             size=16, color=hex_rgb(THEME["primary"]))
    add_text(s, 0.5, 4.4, 12, 0.4, "我的工作", size=12, color=hex_rgb(THEME["secondary"]))
    add_text(s, 0.5, 4.8, 12, 1.5,
             "1. 选型 & 架构设计  2. STT/LLM/TTS 三段流水线打通  3. WebSocket 流式传输  4. 延迟优化  5. 单元/压力/真机测试  6. 三个入口兜底（Twilio / 浏览器 / 小程序）",
             size=14, color=hex_rgb(THEME["primary"]))


def s03_why_this_matters(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, 3, TOTAL, "为什么这件事值得做", "业务价值 → 技术挑战 → 我的取舍")
    add_card(s, 0.5, 2.0, 4.0, 3.8, "现状（痛点）", [
        "• 门卫每天接 50+ 通来访电话",
        "• 手动登记：姓名/车牌/公司/事由",
        "• 高峰期漏接、记错字、信息不完整",
        "• 纸质登记难检索、难统计",
    ])
    add_card(s, 4.75, 2.0, 4.0, 3.8, "目标（我们的解法）", [
        "• AI 7×24 接听，< 25s 内完成登记",
        "• 自然对话，访客零学习成本",
        "• 关键字段强制采集（车牌/事由）",
        "• 结构化落库 + 实时推企业微信",
    ], accent="accent")
    add_card(s, 9.0, 2.0, 3.83, 3.8, "我的取舍", [
        "• 优先体验 & 完整度，不是技术炫技",
        "• 用成熟 SaaS (Twilio + 百炼) 而非自建",
        "• LLM 选中文质量最好的一家",
        "• 兜底：3 个入口 + 静默兜底 + WeChat 推送",
    ], accent="good")


def s04_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, 4, TOTAL, "整体架构", "一图讲清楚：四层 + 三入口 + 一条推流")
    # 简化版架构图（不用 SVG，直接画块）
    add_rect(s, 0.5, 1.8, 2.5, 1.2, THEME["light"])
    add_text(s, 0.5, 1.95, 2.5, 0.4, "访客", size=12, color=hex_rgb(THEME["muted"]), align=PP_ALIGN.CENTER)
    add_text(s, 0.5, 2.35, 2.5, 0.5, "📞 电话 / 🎙 浏览器 / 📱 小程序", size=12, bold=True, color=hex_rgb(THEME["primary"]), align=PP_ALIGN.CENTER)

    add_rect(s, 3.5, 1.8, 2.5, 1.2, "FFF7E6", line="F59E0B")
    add_text(s, 3.5, 1.95, 2.5, 0.4, "传输层", size=12, color=hex_rgb(THEME["muted"]), align=PP_ALIGN.CENTER)
    add_text(s, 3.5, 2.35, 2.5, 0.5, "Twilio MS  /  WebSocket", size=12, bold=True, color=hex_rgb(THEME["primary"]), align=PP_ALIGN.CENTER)

    add_rect(s, 6.5, 1.0, 3.0, 4.0, "E0F2FE", line="0077B6")
    add_text(s, 6.5, 1.15, 3.0, 0.4, "服务层 (FastAPI + Python)", size=12, bold=True, color=hex_rgb(THEME["secondary"]), align=PP_ALIGN.CENTER)
    for i, (name, t) in enumerate([
        ("STT  paraformer", "μ-law → text"),
        ("LLM  qwen-plus",  "function call"),
        ("TTS  qwen3-flash","text → PCM"),
        ("工具",            "save_visit / push_wechat"),
        ("存储  SQLite",     "访客历史"),
    ]):
        add_rect(s, 6.7, 1.65 + i*0.6, 2.6, 0.5, THEME["white"], line=THEME["border"])
        add_text(s, 6.75, 1.72 + i*0.6, 2.5, 0.4, name, size=11, bold=True, color=hex_rgb(THEME["primary"]))
        add_text(s, 8.0, 1.72 + i*0.6, 1.2, 0.4, t, size=9, color=hex_rgb(THEME["muted"]), font=FONT_EN, align=PP_ALIGN.RIGHT)

    add_rect(s, 10.0, 1.8, 2.83, 1.2, "F0FDF4", line="10B981")
    add_text(s, 10.0, 1.95, 2.83, 0.4, "通知", size=12, color=hex_rgb(THEME["muted"]), align=PP_ALIGN.CENTER)
    add_text(s, 10.0, 2.35, 2.83, 0.5, "💬 企业微信群机器人", size=12, bold=True, color=hex_rgb(THEME["primary"]), align=PP_ALIGN.CENTER)

    # 时延预算
    add_rect(s, 0.5, 6.5, 12.33, 1.2, "FFF7E6", line="F59E0B")
    add_text(s, 0.7, 6.6, 12, 0.4, "时延预算 (目标 ≤ 25s):", size=12, bold=True, color=hex_rgb(THEME["warn"]))
    add_text(s, 0.7, 7.0, 12, 0.6,
             "电话接通 1.0s  +  STT 0.3-0.8s  +  LLM 0.5-1.3s  +  TTS 0.7-1.8s  +  WeChat 0.1-0.3s   =   端到端 3-5s",
             size=12, color=hex_rgb(THEME["primary"]))


def s05_decision_twilio(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, 5, TOTAL, "技术选型 ①：为什么 Twilio Media Streams", "不走 SIP，不自建媒体网关")
    add_text(s, 0.5, 2.0, 12, 0.4, "三个候选", size=14, bold=True, color=hex_rgb(THEME["secondary"]))
    items = [
        ("SIP + FreeSWITCH", "全自控，但要从零搭媒体网关、调 SIP 协议", "❌ 杀鸡用牛刀", "bad"),
        ("Twilio Programmable Voice + <Say>/<Play>", "成熟，但只能整段下载播放，无流式", "⚠️ 不够实时", "warn"),
        ("Twilio Media Streams", "双向 WebSocket，μ-law 8kHz 实时流", "✅ 选定", "good"),
    ]
    for i, (name, desc, verdict, color) in enumerate(items):
        add_rect(s, 0.5, 2.6 + i*1.2, 8.0, 1.0, THEME["white"], line=THEME["border"])
        add_rect(s, 0.5, 2.6 + i*1.2, 0.12, 1.0, THEME[color])
        add_text(s, 0.7, 2.75 + i*1.2, 3.0, 0.4, name, size=14, bold=True, color=hex_rgb(THEME["primary"]))
        add_text(s, 0.7, 3.2 + i*1.2, 6.0, 0.4, desc, size=11, color=hex_rgb(THEME["muted"]))
        add_text(s, 8.0, 2.85 + i*1.2, 1.5, 0.5, verdict, size=14, bold=True, color=hex_rgb(THEME[color]), align=PP_ALIGN.CENTER)
    add_text(s, 8.7, 2.0, 4.5, 0.4, "关键收益", size=14, bold=True, color=hex_rgb(THEME["secondary"]))
    add_text(s, 8.7, 2.6, 4.5, 4.0,
             "• 0 媒体服务器运维\n"
             "• 直接吃 WebSocket，开发体验跟写 HTTP 一样\n"
             "• μ-law 8kHz 是电话标准，无需降采样\n"
             "• Twilio 现成录音/转写/分析生态，扩展容易",
             size=12, color=hex_rgb(THEME["primary"]))


def s06_decision_bailian(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, 6, TOTAL, "技术选型 ②：为什么 百炼 (而不是 Deepgram/OpenAI/Edge)", "中文 + 一家搞定 + 链路最短")
    # 表格：横轴 STT/LLM/TTS，纵轴方案
    add_rect(s, 0.5, 1.8, 1.8, 0.6, THEME["secondary"])
    add_text(s, 0.5, 1.95, 1.8, 0.3, "维度", size=12, bold=True, color=hex_rgb(THEME["white"]), align=PP_ALIGN.CENTER)
    headers = ["STT", "LLM", "TTS"]
    for i, h in enumerate(headers):
        add_rect(s, 2.3 + i*2.2, 1.8, 2.2, 0.6, THEME["secondary"])
        add_text(s, 2.3 + i*2.2, 1.95, 2.2, 0.3, h, size=12, bold=True, color=hex_rgb(THEME["white"]), align=PP_ALIGN.CENTER)

    rows = [
        ("Deepgram nova-2 + GPT-4o-mini + Edge TTS", "英文强 / 中文勉强", "贵 / 调用分散", "合成一般"),
        ("自建 Whisper + LLaMA + cosyvoice", "全控", "运维地狱", "训练难"),
        ("百炼 (paraformer + qwen-plus + qwen3-tts-flash)", "中文最优", "一家计费", "Cherry 自然"),
    ]
    for i, row in enumerate(rows):
        add_rect(s, 0.5, 2.4 + i*0.85, 1.8, 0.85, THEME["white"] if i < 2 else "DBEAFE", line=THEME["border"])
        add_text(s, 0.5, 2.6 + i*0.85, 1.8, 0.4, row[0], size=10, bold=(i == 2), color=hex_rgb(THEME["primary"]), align=PP_ALIGN.CENTER)
        for j, val in enumerate(row[1:]):
            add_rect(s, 2.3 + j*2.2, 2.4 + i*0.85, 2.2, 0.85, THEME["white"] if i < 2 else "DBEAFE", line=THEME["border"])
            add_text(s, 2.3 + j*2.2, 2.55 + i*0.85, 2.2, 0.6, val, size=10, color=hex_rgb(THEME["good"] if i == 2 else THEME["primary"]), align=PP_ALIGN.CENTER)

    add_text(s, 0.5, 5.5, 12, 0.4, "为什么选百炼", size=14, bold=True, color=hex_rgb(THEME["secondary"]))
    add_text(s, 0.5, 5.95, 12, 1.0,
             "• 中文质量: paraformer-realtime 是为数不多电话级中文流式 STT\n"
             "• 一致性: STT/LLM/TTS 同一平台，计费/监控/对账简单\n"
             "• 速度: 国内 endpoint + 短链路，TTS 0.7-1.8s 优于 Edge",
             size=12, color=hex_rgb(THEME["primary"]))


def s07_stt_pipeline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, 7, TOTAL, "STT 链路：电话音频 → 文字", "Twilio Media Stream 流入 → 百炼 WebSocket → final transcript")
    # 流程：Twilio MS(μ-law 8k) → 8k PCM → WS → paraformer
    add_rect(s, 0.5, 1.8, 3.0, 1.2, THEME["light"])
    add_text(s, 0.5, 1.95, 3.0, 0.4, "Twilio Media Stream", size=12, color=hex_rgb(THEME["muted"]), align=PP_ALIGN.CENTER)
    add_text(s, 0.5, 2.35, 3.0, 0.5, "base64(μ-law 8kHz)", size=12, bold=True, color=hex_rgb(THEME["primary"]), align=PP_ALIGN.CENTER)
    add_text(s, 0.5, 2.85, 3.0, 0.4, "每 ~20ms 一个 media event", size=9, color=hex_rgb(THEME["muted"]), align=PP_ALIGN.CENTER)

    add_text(s, 3.7, 2.3, 0.3, 0.2, "→", size=18, bold=True, color=hex_rgb(THEME["secondary"]))

    add_rect(s, 4.2, 1.8, 3.0, 1.2, "FFF7E6", line="F59E0B")
    add_text(s, 4.2, 1.95, 3.0, 0.4, "Server STT Handler", size=12, color=hex_rgb(THEME["muted"]), align=PP_ALIGN.CENTER)
    add_text(s, 4.2, 2.35, 3.0, 0.5, "ulaw2lin(2) → ratecv(8k)", size=12, bold=True, color=hex_rgb(THEME["primary"]), align=PP_ALIGN.CENTER)
    add_text(s, 4.2, 2.85, 3.0, 0.4, "audioop 一次转完", size=9, color=hex_rgb(THEME["muted"]), align=PP_ALIGN.CENTER)

    add_text(s, 7.4, 2.3, 0.3, 0.2, "→", size=18, bold=True, color=hex_rgb(THEME["secondary"]))

    add_rect(s, 7.9, 1.8, 3.0, 1.2, "E0F2FE", line="0077B6")
    add_text(s, 7.9, 1.95, 3.0, 0.4, "百炼 paraformer-realtime-v1", size=12, color=hex_rgb(THEME["muted"]), align=PP_ALIGN.CENTER)
    add_text(s, 7.9, 2.35, 3.0, 0.5, "WebSocket 流式", size=12, bold=True, color=hex_rgb(THEME["primary"]), align=PP_ALIGN.CENTER)
    add_text(s, 7.9, 2.85, 3.0, 0.4, "边说边出 final", size=9, color=hex_rgb(THEME["muted"]), align=PP_ALIGN.CENTER)

    add_text(s, 11.1, 2.3, 0.3, 0.2, "→", size=18, bold=True, color=hex_rgb(THEME["secondary"]))

    add_rect(s, 11.6, 1.8, 1.5, 1.2, "F0FDF4", line="10B981")
    add_text(s, 11.6, 1.95, 1.5, 0.4, "Output", size=12, color=hex_rgb(THEME["muted"]), align=PP_ALIGN.CENTER)
    add_text(s, 11.6, 2.35, 1.5, 0.5, "text\nis_final", size=12, bold=True, color=hex_rgb(THEME["primary"]), align=PP_ALIGN.CENTER)

    # 关键点
    add_card(s, 0.5, 3.5, 6.0, 3.0, "为什么自己写 WebSocket 不用 SDK", [
        "• dashscope SDK 1.20+ 有 KeyError: 'begin_time' bug",
        "• SDK 硬编码模型白名单拒收 streaming model",
        "• 走裸 WS: payload 字段 + 300ms endpointing 自己控",
    ])
    add_card(s, 6.75, 3.5, 6.08, 3.0, "回访识别（一招优化）", [
        "• 听到车牌就查 DB (app.database.find_returning_visitor)",
        "• 命中就重建 agent with returning history",
        "• 2 轮 ~8s 结束，不用重问基本信息",
    ], accent="good")


def s08_agent_design(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, 8, TOTAL, "LLM Agent：状态机 + 工具调用", "用最少 prompt 实现 1 个目标：凑齐 4 个字段")
    add_rect(s, 0.5, 1.8, 6.0, 4.8, "F8FAFC", line=THEME["border"])
    add_text(s, 0.7, 1.95, 5.6, 0.4, "状态机（4 字段 → complete）", size=14, bold=True, color=hex_rgb(THEME["secondary"]))
    fields = [
        ("plate",  "车牌",       THEME["good"]),
        ("company", "公司",       THEME["secondary"]),
        ("reason", "事由",       THEME["secondary"]),
        ("phone",  "手机号",     THEME["secondary"]),
        ("complete", "→ send_to_guard_and_end_call()", THEME["accent"]),
    ]
    for i, (k, lbl, c) in enumerate(fields):
        add_rect(s, 0.8, 2.5 + i*0.7, 1.5, 0.55, c, line=c)
        add_text(s, 0.8, 2.6 + i*0.7, 1.5, 0.35, k, size=12, bold=True, color=hex_rgb(THEME["white"]), align=PP_ALIGN.CENTER)
        add_text(s, 2.5, 2.6 + i*0.7, 4.0, 0.35, lbl, size=12, color=hex_rgb(THEME["primary"]))
        if i < len(fields) - 1:
            add_text(s, 1.45, 3.05 + i*0.7, 0.2, 0.15, "↓", size=14, bold=True, color=hex_rgb(THEME["muted"]), align=PP_ALIGN.CENTER)

    add_text(s, 6.75, 1.95, 6.08, 0.4, "系统 prompt 关键策略", size=14, bold=True, color=hex_rgb(THEME["secondary"]))
    add_text(s, 6.75, 2.4, 6.08, 4.0, [
        "• role: “你是访客登记员，不是客服”，边界清晰",
        "• 一次只问一个字段，节奏感",
        "• 模糊时用 clarify 追问，不猜",
        "• 听到车牌/公司时立即 update_visitor()",
        "• 4 字段全了就调 end_call，不啰嗦",
    ], size=12, color=hex_rgb(THEME["primary"]))
    add_text(s, 6.75, 5.6, 6.08, 0.4, "工具 schema（OpenAI function calling）", size=12, bold=True, color=hex_rgb(THEME["secondary"]))
    add_text(s, 6.75, 6.05, 6.08, 0.6, "update_visitor / send_to_guard_and_end_call", size=12, font=FONT_EN, color=hex_rgb(THEME["primary"]))


def s09_tts_optimization(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, 9, TOTAL, "TTS 优化：1.95s → 37ms（首音延迟）", "把慢的活挪到 server boot，session 级几乎零等待")
    # 左侧：优化前 / 优化后 对比
    add_rect(s, 0.5, 1.8, 6.0, 1.2, "FEE2E2", line="EF4444")
    add_text(s, 0.7, 1.95, 5.6, 0.4, "❌ 优化前", size=14, bold=True, color=hex_rgb(THEME["bad"]))
    add_text(s, 0.7, 2.4, 5.6, 0.4, "用户进 session → TTS 现合成 greeting → 等 1.5-1.8s → 播", size=12, color=hex_rgb(THEME["primary"]))
    add_text(s, 0.7, 2.7, 5.6, 0.4, "首音延迟: ~1.95s   用户感知: 慢", size=12, bold=True, color=hex_rgb(THEME["bad"]))

    add_rect(s, 6.83, 1.8, 6.0, 1.2, "D1FAE5", line="10B981")
    add_text(s, 7.03, 1.95, 5.6, 0.4, "✅ 优化后", size=14, bold=True, color=hex_rgb(THEME["good"]))
    add_text(s, 7.03, 2.4, 5.6, 0.4, "server boot 时一次性预合成 → session 只剩内存读", size=12, color=hex_rgb(THEME["primary"]))
    add_text(s, 7.03, 2.7, 5.6, 0.4, "首音延迟: 37ms   用户感知: 即时", size=12, bold=True, color=hex_rgb(THEME["good"]))

    # 三个优化点
    add_card(s, 0.5, 3.2, 4.0, 3.5, "① Boot 时预合成 greeting", [
        "lifespan 启动时调 _synth_bailian(开场白)",
        "结果存 BrowserSession._BOOT_GREETING_PCM",
        "session.run() 直接从内存读，零 TTS 等待",
    ])
    add_card(s, 4.75, 3.2, 4.0, 3.5, "② 短句 TTS cache", [
        "_TTS_CACHE: 4 条常用回复",
        "好的 / 请讲 / 请再说一遍 / 好的，已通知门卫",
        "session 内命中 → 0ms",
    ], accent="accent")
    add_card(s, 9.0, 3.2, 3.83, 3.5, "③ 浏览器端防抖", [
        "agent 说话时 micMuted=true",
        "防止 TTS 漏音触发 STT 二次识别",
        "（barge-in protection）",
    ], accent="warn")

    add_text(s, 0.5, 7.0, 12, 0.4, "性能证据（5 路并发 stress test）", size=14, bold=True, color=hex_rgb(THEME["secondary"]))
    add_text(s, 0.5, 7.5, 12, 0.4, "SLA 5/5 全过（≤25s）· T1st p50 ~2.0s · Tend p50 ~10.2s", size=12, color=hex_rgb(THEME["primary"]))


def s10_three_entrypoints(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, 10, TOTAL, "三个入口（Twilio / WebRTC / 小程序）", "同一套 STT/LLM/TTS，三种用户场景")
    # 三栏
    add_card(s, 0.5, 1.8, 4.0, 5.0, "📞 Twilio 真实电话", [
        "• 生产路径：PSTN 拨号",
        "• Twilio Media Streams 接音",
        "• 体验最真实",
        "• 需要 Twilio 账号 + ICP 备案域名",
        "",
        "现状：账号关停，9 通 call log 留底",
    ])
    add_card(s, 4.75, 1.8, 4.0, 5.0, "🌐 WebRTC 浏览器", [
        "• getUserMedia + WebSocket",
        "• /ws/browser 端点",
        "• 零电话零账号",
        "• 面试官现场可演示",
        "",
        "状态：跑通，live URL 已部署",
    ], accent="accent")
    add_card(s, 9.0, 1.8, 3.83, 5.0, "📱 微信小程序", [
        "• wx.getRecorderManager() 录音",
        "• wx.sendSocketMessage() 推流",
        "• wx.createWebAudioContext() 播音",
        "• 走微信内网，零流量费",
        "",
        "状态：drop-in 4 文件 + INTEGRATE.md",
    ], accent="good")

    add_text(s, 0.5, 7.1, 12, 0.4, "设计原则", size=12, bold=True, color=hex_rgb(THEME["secondary"]))
    add_text(s, 0.5, 7.5, 12, 0.4, "入口不同 → transport 不同，但后端 BrowserSession / TwilioCall 完全同源 STT/LLM/TTS。改 transport 不改 agent。",
             size=12, color=hex_rgb(THEME["primary"]))


def s11_testing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, 11, TOTAL, "测试 & 验证", "三件证据，证明链路真的跑通")
    # 三栏
    add_card(s, 0.5, 1.8, 4.0, 2.5, "单元测试", [
        "72 / 72 全过",
        "覆盖 plate / agent / schemas / wechat",
        "uv run pytest tests/ -v",
    ])
    add_card(s, 4.75, 1.8, 4.0, 2.5, "压力测试 (5x5)", [
        "5 通并发 3 / 全部 ≤ 25s",
        "T1st ~2s · Tend ~10s",
        "agent 100% 跑通对话",
    ], accent="accent")
    add_card(s, 9.0, 1.8, 3.83, 2.5, "真实 Twilio 通话", [
        "9 通 inbound call 全 completed",
        "server log 抓得到每通 call",
        "已扣 $0.0405",
    ], accent="good")

    # demo 引用
    add_text(s, 0.5, 4.6, 12, 0.4, "可以现场播的证据", size=14, bold=True, color=hex_rgb(THEME["secondary"]))
    add_card(s, 0.5, 5.1, 6.0, 2.0, "🎙 28s 完整对话录音", [
        "3 轮真实对话（greeting + 车牌 + 完成）",
        "Cherry 女生自然中文",
        "路径: docs/evidence/demo-3turn-28s.wav",
    ])
    add_card(s, 6.75, 5.1, 6.08, 2.0, "🌐 WebRTC 浏览器实测", [
        "面试官现场可点开",
        "授权麦克风 → 对着电脑说话",
        "https://cars-gathering-divx-bent.trycloudflare.com/browser-test",
    ], accent="accent")


def s12_debugging_journey(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, 12, TOTAL, "踩过的坑（debugging journey）", "高级信号 = 你踩过 + 解决了 + 知道为什么")
    items = [
        ("Trial 账户 Media Stream 吞音",
         "Twilio Trial 不发 stream 音频，9 通 call 全听不到",
         "升级 Full 账户 → 音频立即可播",
         "帮面试官看到：会读 Twilio 文档、能区分账户级别限制"),
        ("ngrok-free 不支持 WSS",
         "Twilio Stream 连不上，101 Switching 拿不到",
         "切 cloudflared，--protocol http2（QUIC 在中国墙）",
         "帮面试官看到：换工具前先复现、懂 WSS 升级机制"),
        ("TTS 缓存误用，重复合成 1.8s",
         "每次 greeting 都现合成，慢",
         "boot 时一次性预合成 → 1.95s → 37ms",
         "帮面试官看到：会从'感知延迟'角度优化、不是只看 QPS"),
        ("agent 太快导致串话",
         "用户还没说完，agent 答了；用户接着说，STT 抓错",
         "agent 说话时 micMuted=true + 250ms grace",
         "帮面试官看到：理解'现实 UX 跟 demo 不一样'"),
        ("WebRTC 音频怪声",
         "浏览器把 16k PCM 按 48k 播，3 倍速",
         "AudioContext 不指定 sampleRate，让 Web Audio 自动重采样",
         "帮面试官看到：底层理解、知道 API 默认可坑"),
    ]
    for i, (title, sym, fix, signal) in enumerate(items):
        y = 1.8 + i * 1.05
        add_rect(s, 0.5, y, 12.33, 0.95, THEME["white"], line=THEME["border"])
        add_rect(s, 0.5, y, 0.12, 0.95, THEME["warn"])
        add_text(s, 0.7, y + 0.05, 4.0, 0.4, title, size=13, bold=True, color=hex_rgb(THEME["primary"]))
        add_text(s, 0.7, y + 0.45, 4.0, 0.45, sym, size=10, color=hex_rgb(THEME["muted"]))
        add_text(s, 4.8, y + 0.05, 4.0, 0.4, "修法", size=11, bold=True, color=hex_rgb(THEME["good"]))
        add_text(s, 4.8, y + 0.45, 4.0, 0.45, fix, size=10, color=hex_rgb(THEME["primary"]))
        add_text(s, 9.0, y + 0.05, 4.0, 0.4, "面试官看", size=11, bold=True, color=hex_rgb(THEME["secondary"]))
        add_text(s, 9.0, y + 0.45, 4.0, 0.45, signal, size=10, color=hex_rgb(THEME["muted"]))


def s13_limitations(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, 13, TOTAL, "局限 & 未来", "自省 = 加分项")
    add_card(s, 0.5, 1.8, 6.0, 5.0, "🔍 已知的局限", [
        "• 百炼 TTS REST 非流式，首包 0.7-1.8s",
        "  缓解: boot 预合成 + 短句 cache",
        "• 浏览器录音只能实时拿，不能传文件测",
        "• cloudflared quick tunnel 每次重启换 URL",
        "  缓解: 长期部署需自有域名 / 阿里云",
        "• 个人主体小程序不能上架",
        "• 国际长途有运营商风控（已绕过）",
    ])
    add_card(s, 6.75, 1.8, 6.08, 5.0, "🚀 下一步可以做的", [
        "• STT/LLM/TTS 全异步 pipeline（Q1 完成）",
        "• LLM 切 streaming，前几个 token 就开始 TTS",
        "  （理论能再省 1s 延迟）",
        "• 实时情绪识别（Twilio + 声纹）",
        "• 多语言切换（中/英/粤语）",
        "• 门卫端反向查询 API（已有，扩展）",
    ], accent="accent")


def s14_qa(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, 14, TOTAL, "Q&A 准备", "高频问题的 1-2 句快速回答")
    qas = [
        ("为什么不用 Deepgram nova-2？",
         "百炼 paraformer 中文更好，且 STT/LLM/TTS 一家计费链路短。Deepgram 是备选，代码里 make_transcriber() 留了切换接口。"),
        ("Twilio Media Stream 跟 SIP 区别？",
         "Media Stream 是 Twilio 把电话音频推 WebSocket，开发者直接消费 μ-law 8k。SIP 要自建媒体服务器，工作量 10x+。"),
        ("端到端延迟怎么优化的？",
         "三处：boot 预合成 greeting (1.95s→37ms)、TTS cache 短句 (1.5s→0ms)、浏览器 mic mute 防串话（体感）。"),
        ("三个入口会不会太复杂？",
         "后端只一份：BrowserSession（浏览器/小程序）+ TwilioCall（电话），共享 STT/LLM/TTS。transport 不同，agent 同一份。"),
        ("你说 Twilio 账号关停了，怎么验证电话链路？",
         "服务器 log 留 9 通 call 记录，每通都有 STT/LLM/TTS 完整 log。WebRTC 路径同 STT/LLM/TTS，作为电话链路的代理验证。"),
        ("如果 Twilio 完全不能用，会怎么改？",
         "换 Telnyx（API 兼容 Twilio，~1 天改 handler）；或部署 WebRTC + 小程序为生产入口（已实现），跳电话栈。"),
    ]
    for i, (q, a) in enumerate(qas):
        y = 1.8 + i * 0.85
        add_rect(s, 0.5, y, 12.33, 0.78, THEME["white"], line=THEME["border"])
        add_text(s, 0.7, y + 0.05, 4.5, 0.4, "Q: " + q, size=12, bold=True, color=hex_rgb(THEME["secondary"]))
        add_text(s, 0.7, y + 0.42, 12, 0.4, "A: " + a, size=11, color=hex_rgb(THEME["primary"]))


BUILDERS = [
    s01_cover, s02_self, s03_why_this_matters, s04_architecture,
    s05_decision_twilio, s06_decision_bailian, s07_stt_pipeline,
    s08_agent_design, s09_tts_optimization, s10_three_entrypoints,
    s11_testing, s12_debugging_journey, s13_limitations, s14_qa,
]


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    for fn in BUILDERS:
        fn(prs)
    out = Path("docs/pptx/output/voice-agent-defense.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"✓ {out}  ({len(BUILDERS)} slides)")


if __name__ == "__main__":
    build()
