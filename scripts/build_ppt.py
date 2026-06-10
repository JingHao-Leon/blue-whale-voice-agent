"""
Build a professional PPT for the Voice Agent take-home submission.

Theme: Pure Tech Blue (palette #15 from the design-system).
Style: Soft & Balanced.

Run: python3 scripts/build_ppt.py
Output: docs/pptx/output/voice-agent-take-home.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from lxml import etree


# ---------------------------------------------------------------------------
# Theme — Pure Tech Blue
# ---------------------------------------------------------------------------
# Source palette: 03045e 0077b6 00b4d8 90e0ef caf0f8
THEME = {
    "primary":   "03045E",   # deep navy — titles, body
    "secondary": "0077B6",   # medium blue — accents, callouts
    "accent":    "00B4D8",   # bright cyan — highlights
    "light":     "CAF0F8",   # very light blue — card backgrounds
    "bg":        "F8FAFC",   # very light gray — slide background
    "muted":     "64748B",   # gray-500 — secondary text
    "border":    "E2E8F0",   # gray-200 — borders
    "white":     "FFFFFF",
    "good":      "10B981",   # green for "good" callout
    "bad":       "EF4444",   # red for "bad" callout
    "warn":      "F59E0B",   # amber for warn / SLA
}

FONT_CN = "Microsoft YaHei"
FONT_EN = "Arial"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def hex_rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str)


def add_text(
    slide,
    x: float, y: float, w: float, h: float,
    text: str | list,
    *,
    font_size: int = 14,
    bold: bool = False,
    color: str = "primary",
    font: str = FONT_CN,
    align: str = "left",
    valign: str = "top",
    line_spacing: float | None = None,
):
    """Add a text box. `text` may be a string or list of (text, options) tuples."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    if isinstance(text, str):
        text = [(text, {})]

    for i, (t, opts) in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = t
        run.font.name = opts.get("font", font)
        run.font.size = Pt(opts.get("size", font_size))
        run.font.bold = opts.get("bold", bold)
        c = opts.get("color", color)
        run.font.color.rgb = hex_rgb(THEME[c] if c in THEME else c)
    return tb


def add_rect(slide, x: float, y: float, w: float, h: float, *, fill: str, line: str | None = None, line_w: float = 0.75, radius: float | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE if radius is None else MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(THEME[fill] if fill in THEME else fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hex_rgb(THEME[line] if line in THEME else line)
        shape.line.width = Pt(line_w)
    if radius is not None:
        shape.adjustments[0] = radius
    shape.shadow.inherit = False
    return shape


def add_line(slide, x1: float, y1: float, x2: float, y2: float, *, color: str = "muted", width: float = 1.0):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = hex_rgb(THEME[color] if color in THEME else color)
    line.line.width = Pt(width)
    return line


def add_page_number(slide, n: int, total: int = 10):
    add_text(
        slide, 9.0, 5.25, 0.9, 0.3,
        f"{n:02d} / {total:02d}",
        font_size=9, color="muted", font=FONT_EN, align="right",
    )


def set_bg(slide, color: str = "bg"):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_rgb(THEME[color])


def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 6 = blank


# ---------------------------------------------------------------------------
# Build slides
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    sw, sh = prs.slide_width, prs.slide_height
    total = 10

    # =================== Slide 01 — Cover ===================
    s = slide_blank(prs)
    set_bg(s, "bg")

    # Left accent block
    add_rect(s, 0, 0, 0.5, 5.625, fill="secondary")

    # Big title
    add_text(
        s, 1.0, 1.2, 8.5, 1.0,
        [("Voice Agent", {"size": 56, "bold": True, "color": "primary", "font": FONT_EN})],
    )
    add_text(
        s, 1.0, 2.2, 8.5, 0.8,
        [("\u84dd\u8272\u9cb8\u9c7c\u79d1\u6280\u56ed \u00b7 \u8bbf\u5ba2\u767b\u8bb0\u7cfb\u7edf", {"size": 32, "bold": True, "color": "secondary", "font": FONT_CN})],
    )

    # Subtitle
    add_text(
        s, 1.0, 3.1, 8.5, 0.4,
        [("Take-Home Exam Submission \u00b7 2026", {"size": 14, "color": "muted", "font": FONT_EN})],
    )

    # Tag pills
    add_rect(s, 1.0, 3.7, 1.5, 0.35, fill="light", radius=0.4)
    add_text(s, 1.0, 3.7, 1.5, 0.35, "FastAPI", font_size=10, color="secondary", font=FONT_EN, align="center", valign="middle", bold=True)
    add_rect(s, 2.6, 3.7, 1.5, 0.35, fill="light", radius=0.4)
    add_text(s, 2.6, 3.7, 1.5, 0.35, "Twilio", font_size=10, color="secondary", font=FONT_EN, align="center", valign="middle", bold=True)
    add_rect(s, 4.2, 3.7, 1.7, 0.35, fill="light", radius=0.4)
    add_text(s, 4.2, 3.7, 1.7, 0.35, "Deepgram", font_size=10, color="secondary", font=FONT_EN, align="center", valign="middle", bold=True)
    add_rect(s, 6.0, 3.7, 1.5, 0.35, fill="light", radius=0.4)
    add_text(s, 6.0, 3.7, 1.5, 0.35, "GPT-4o", font_size=10, color="secondary", font=FONT_EN, align="center", valign="middle", bold=True)
    add_rect(s, 7.6, 3.7, 1.4, 0.35, fill="light", radius=0.4)
    add_text(s, 7.6, 3.7, 1.4, 0.35, "Edge TTS", font_size=10, color="secondary", font=FONT_EN, align="center", valign="middle", bold=True)

    # Bottom credit
    add_text(
        s, 1.0, 4.7, 8.5, 0.4,
        "Whale Tech Take-Home Exam \u00b7 Voice Agent for Industrial Park",
        font_size=12, color="muted", font=FONT_EN,
    )
    add_text(
        s, 1.0, 5.05, 8.5, 0.3,
        "\u4e0a\u6d77\u84dd\u8272\u9cb8\u9c7c\u79d1\u6280\u6709\u9650\u516c\u53f8",
        font_size=11, color="muted", font=FONT_CN,
    )

    # =================== Slide 02 — TOC ===================
    s = slide_blank(prs)
    set_bg(s, "bg")
    add_text(s, 0.6, 0.5, 9, 0.7, "Agenda", font_size=36, bold=True, color="primary", font=FONT_EN)
    add_text(s, 0.6, 1.1, 9, 0.4, "\u62a5\u544a\u76ee\u5f55", font_size=16, color="muted", font=FONT_CN)
    add_line(s, 0.6, 1.6, 9.4, 1.6, color="border", width=1.0)

    sections = [
        ("01", "\u4e1a\u52a1\u80cc\u666f\u4e0e\u75db\u70b9", "Industrial Park Visitor Pain Points"),
        ("02", "\u7cfb\u7edf\u67b6\u6784\u4e0e\u6280\u672f\u9009\u578b", "System Architecture & Tech Choices"),
        ("03", "\u5bf9\u8bdd\u4f53\u9a8c\u6807\u51c6", "Conversation Quality Standards"),
        ("04", "25 \u79d2 SLA \u65f6\u5ef6\u62c6\u89e3", "End-to-End Latency Breakdown"),
        ("05", "\u4f01\u4e1a\u5fae\u4fe1\u96c6\u6210\u4e0e\u4ea4\u4ed8", "WeChat Work Integration & Delivery"),
    ]
    y = 1.9
    for num, cn, en in sections:
        add_text(s, 0.7, y, 0.8, 0.5, num, font_size=28, bold=True, color="accent", font=FONT_EN, valign="middle")
        add_text(s, 1.6, y, 4.2, 0.5, cn, font_size=20, bold=True, color="primary", font=FONT_CN, valign="middle")
        add_text(s, 5.9, y, 4.0, 0.5, en, font_size=12, color="muted", font=FONT_EN, valign="middle")
        y += 0.6

    add_page_number(s, 2, total)

    # =================== Slide 03 — Section 01 Divider ===================
    s = slide_blank(prs)
    set_bg(s, "bg")
    # Left block big number
    add_rect(s, 0.6, 0.6, 4.5, 4.4, fill="light", radius=0.1)
    add_text(s, 0.6, 1.1, 4.5, 2.5, "01", font_size=200, bold=True, color="secondary", font=FONT_EN, align="center", valign="middle")
    add_text(s, 0.6, 3.6, 4.5, 0.4, "Section", font_size=14, color="muted", font=FONT_EN, align="center")
    add_text(s, 0.6, 4.0, 4.5, 0.4, "PART ONE", font_size=12, color="muted", font=FONT_EN, align="center", bold=True)

    add_text(s, 5.4, 1.5, 4.2, 0.6, "\u4e1a\u52a1\u80cc\u666f\u4e0e\u75db\u70b9", font_size=36, bold=True, color="primary", font=FONT_CN)
    add_text(s, 5.4, 2.3, 4.2, 0.4, "Business Context & Pain Points", font_size=14, color="secondary", font=FONT_EN)
    add_line(s, 5.4, 2.85, 9.4, 2.85, color="accent", width=2.0)
    add_text(
        s, 5.4, 3.05, 4.2, 1.5,
        "\u4ece\u4eba\u5de5\u95ee\u8be2\u5230 AI \u63a5\u542c\uff1a\n\u89e3\u51b3\u56ed\u533a\u8bbf\u5ba2\u767b\u8bb0\u4f4e\u6548\u3001\u4eba\u529b\u6d6a\u8d39\u7684\u95ee\u9898",
        font_size=14, color="muted", font=FONT_CN, line_spacing=1.5,
    )

    add_page_number(s, 3, total)

    # =================== Slide 04 — Business context ===================
    s = slide_blank(prs)
    set_bg(s, "bg")
    add_text(s, 0.6, 0.4, 9, 0.6, "\u4e1a\u52a1\u573a\u666f\uff1a\u56ed\u533a\u8bbf\u5ba2\u767b\u8bb0", font_size=28, bold=True, color="primary", font=FONT_CN)
    add_text(s, 0.6, 0.95, 9, 0.4, "Industrial Park Visitor Registration", font_size=12, color="muted", font=FONT_EN)
    add_line(s, 0.6, 1.35, 9.4, 1.35, color="border", width=1.0)

    # Left: Pain point
    add_text(s, 0.6, 1.6, 4.4, 0.4, "\u2460 \u5f53\u524d\u6d41\u7a0b\uff08\u4eba\u5de5\uff09", font_size=16, bold=True, color="primary", font=FONT_CN)
    flow_steps = [
        "1. \u8bbf\u5ba2\u8f66\u5230\u8fbe\u5165\u53e3",
        "2. \u6d77\u5eb7\u95e8\u7981\u8bc6\u522b\u8f66\u724c\uff1a\u201c\u672a\u767b\u8bb0\u201d",
        "3. \u4fdd\u5b89\u5e26\u7eb8\u8d28\u767b\u8bb0\u8868\u8d70\u51fa\u4fdd\u5b89\u5ba4",
        "4. \u53e3\u5934\u8be2\u95ee\u5e76\u624b\u5de5\u586b\u5199\u4fe1\u606f",
        "5. \u9065\u63a7\u5668\u624b\u52a8\u62ac\u6746\u653e\u884c",
    ]
    y = 2.05
    for step in flow_steps:
        add_rect(s, 0.6, y, 0.1, 0.35, fill="muted")
        add_text(s, 0.85, y, 4.0, 0.35, step, font_size=12, color="primary", font=FONT_CN, valign="middle")
        y += 0.4

    # Pain callout
    add_rect(s, 0.6, 4.15, 4.4, 1.0, fill="FEF2F2", line="EF4444", radius=0.1)
    add_text(s, 0.75, 4.25, 4.1, 0.3, "\u26a0 \u75db\u70b9", font_size=11, bold=True, color="EF4444", font=FONT_CN)
    add_text(
        s, 0.75, 4.55, 4.1, 0.6,
        "\u4f4e\u6548 \u00b7 \u4eba\u529b\u6d6a\u8d39 \u00b7 \u4f9d\u8d56\u7eb8\u8d28 \u00b7 \u51ac\u5929\u96be\u53d7 \u00b7 \u96be\u4ee5\u5ba1\u8ba1",
        font_size=11, color="primary", font=FONT_CN, line_spacing=1.4,
    )

    # Right: Target flow
    add_text(s, 5.2, 1.6, 4.4, 0.4, "\u2461 \u76ee\u6807\u6d41\u7a0b\uff08AI \u81ea\u52a8\u5316\uff09", font_size=16, bold=True, color="secondary", font=FONT_CN)
    new_steps = [
        "1. \u8bbf\u5ba2\u62e8\u6253\u56ed\u533a\u70ed\u7ebf\u7535\u8bdd",
        "2. AI \u63a5\u542c \u00b7 \u81ea\u7136\u5bf9\u8bdd\u91c7\u96c6\u4fe1\u606f",
        "3. \u7ed3\u6784\u5316\u4fe1\u606f\u63a8\u9001\u81f3\u95e8\u536b\u4f01\u4e1a\u5fae\u4fe1",
        "4. \u95e8\u536b\u786e\u8ba4 \u00b7 \u8fdc\u7a0b\u63a5\u653e\u884c",
    ]
    y = 2.05
    for step in new_steps:
        add_rect(s, 5.2, y, 0.1, 0.35, fill="10B981")
        add_text(s, 5.45, y, 4.0, 0.35, step, font_size=12, color="primary", font=FONT_CN, valign="middle")
        y += 0.4

    # Win callout
    add_rect(s, 5.2, 4.15, 4.4, 1.0, fill="ECFDF5", line="10B981", radius=0.1)
    add_text(s, 5.35, 4.25, 4.1, 0.3, "\u2713 \u4ef7\u503c", font_size=11, bold=True, color="10B981", font=FONT_CN)
    add_text(
        s, 5.35, 4.55, 4.1, 0.6,
        "25 \u79d2\u5b8c\u6210 \u00b7 7\u00d724h \u968f\u65f6 \u00b7 \u53ef\u5ba1\u8ba1 \u00b7 \u591a\u8f66\u5e76\u53d1",
        font_size=11, color="primary", font=FONT_CN, line_spacing=1.4,
    )

    add_page_number(s, 4, total)

    # =================== Slide 05 — Section 02 Divider ===================
    s = slide_blank(prs)
    set_bg(s, "bg")
    add_rect(s, 0.6, 0.6, 4.5, 4.4, fill="light", radius=0.1)
    add_text(s, 0.6, 1.1, 4.5, 2.5, "02", font_size=200, bold=True, color="secondary", font=FONT_EN, align="center", valign="middle")
    add_text(s, 0.6, 3.6, 4.5, 0.4, "Section", font_size=14, color="muted", font=FONT_EN, align="center")
    add_text(s, 0.6, 4.0, 4.5, 0.4, "PART TWO", font_size=12, color="muted", font=FONT_EN, align="center", bold=True)

    add_text(s, 5.4, 1.5, 4.2, 0.6, "\u7cfb\u7edf\u67b6\u6784\u4e0e\u6280\u672f\u9009\u578b", font_size=36, bold=True, color="primary", font=FONT_CN)
    add_text(s, 5.4, 2.3, 4.2, 0.4, "System Architecture & Tech Choices", font_size=14, color="secondary", font=FONT_EN)
    add_line(s, 5.4, 2.85, 9.4, 2.85, color="accent", width=2.0)
    add_text(
        s, 5.4, 3.05, 4.2, 1.5,
        "Twilio \u00b7 Deepgram \u00b7 GPT-4o-mini \u00b7 Edge TTS \u00b7 \u4f01\u4e1a\u5fae\u4fe1 Webhook",
        font_size=14, color="muted", font=FONT_CN, line_spacing=1.5,
    )

    add_page_number(s, 5, total)

    # =================== Slide 06 — Architecture ===================
    s = slide_blank(prs)
    set_bg(s, "bg")
    add_text(s, 0.6, 0.4, 9, 0.6, "\u7cfb\u7edf\u67b6\u6784\u56fe", font_size=28, bold=True, color="primary", font=FONT_CN)
    add_text(s, 0.6, 0.95, 9, 0.4, "End-to-End Architecture \u00b7 Twilio Media Streams + FastAPI", font_size=12, color="muted", font=FONT_EN)
    add_line(s, 0.6, 1.35, 9.4, 1.35, color="border", width=1.0)

    # Driver
    add_rect(s, 0.6, 1.6, 1.2, 0.8, fill="light", radius=0.1)
    add_text(s, 0.6, 1.6, 1.2, 0.4, "🚗", font_size=22, color="primary", align="center", valign="middle")
    add_text(s, 0.6, 2.0, 1.2, 0.4, "\u8bbf\u5ba2\u8f66\u4e3b", font_size=10, bold=True, color="primary", font=FONT_CN, align="center", valign="middle")

    # Twilio
    add_rect(s, 2.0, 1.6, 1.3, 0.8, fill="FEF3C7", radius=0.1)
    add_text(s, 2.0, 1.6, 1.3, 0.3, "\u2601 Twilio", font_size=11, bold=True, color="92400E", font=FONT_EN, align="center", valign="middle")
    add_text(s, 2.0, 1.95, 1.3, 0.4, "Media Streams", font_size=9, color="92400E", font=FONT_EN, align="center", valign="middle")

    # Server container
    add_rect(s, 3.5, 1.5, 3.6, 3.3, fill="white", line="94A3B8", radius=0.1)
    add_rect(s, 3.5, 1.5, 3.6, 0.3, fill="1E293B", radius=0.1)
    add_text(s, 3.5, 1.5, 3.6, 0.3, "\u2699 FastAPI Server", font_size=10, bold=True, color="white", font=FONT_EN, align="center", valign="middle")
    # STT
    add_rect(s, 3.65, 1.95, 1.65, 0.55, fill="E0F2FE", radius=0.05)
    add_text(s, 3.65, 1.95, 1.65, 0.55, "🎙 STT\nDeepgram", font_size=9, bold=True, color="1E40AF", font=FONT_EN, align="center", valign="middle")
    # LLM
    add_rect(s, 5.3, 1.95, 1.65, 0.55, fill="FCE7F3", radius=0.05)
    add_text(s, 5.3, 1.95, 1.65, 0.55, "🧠 LLM Agent\ngpt-4o-mini", font_size=9, bold=True, color="9F1239", font=FONT_EN, align="center", valign="middle")
    # TTS
    add_rect(s, 3.65, 2.6, 1.65, 0.55, fill="E0F2FE", radius=0.05)
    add_text(s, 3.65, 2.6, 1.65, 0.55, "🔊 TTS\nEdge TTS", font_size=9, bold=True, color="1E40AF", font=FONT_EN, align="center", valign="middle")
    # Tools / DB
    add_rect(s, 5.3, 2.6, 1.65, 0.55, fill="F1F5F9", radius=0.05)
    add_text(s, 5.3, 2.6, 1.65, 0.55, "🗜 SQLite\n+ Tools", font_size=9, bold=True, color="0F172A", font=FONT_EN, align="center", valign="middle")
    # WeChat client
    add_rect(s, 3.65, 3.25, 3.3, 0.55, fill="DCFCE7", radius=0.05)
    add_text(s, 3.65, 3.25, 3.3, 0.55, "💬 WeChat Webhook Client (markdown card)", font_size=9, bold=True, color="14532D", font=FONT_EN, align="center", valign="middle")
    # SLA bar
    add_rect(s, 3.65, 3.9, 3.3, 0.7, fill="FEF3C7", radius=0.05)
    add_text(s, 3.65, 3.95, 3.3, 0.3, "\u23F1 SLA \u2264 25 s", font_size=10, bold=True, color="92400E", font=FONT_CN, align="center", valign="middle")
    add_text(s, 3.65, 4.25, 3.3, 0.3, "\u5b9e\u6d4b 3-6 s", font_size=9, color="92400E", font=FONT_CN, align="center", valign="middle")

    # WeChat group
    add_rect(s, 7.3, 1.6, 1.5, 0.8, fill="DCFCE7", radius=0.1)
    add_text(s, 7.3, 1.6, 1.5, 0.3, "💬 \u4f01\u4e1a\u5fae\u4fe1", font_size=11, bold=True, color="14532D", font=FONT_CN, align="center", valign="middle")
    add_text(s, 7.3, 1.95, 1.5, 0.4, "\u95e8\u536b\u7fa4", font_size=9, color="14532D", font=FONT_CN, align="center", valign="middle")

    # Guard
    add_rect(s, 7.3, 2.55, 1.5, 0.7, fill="FCE7F3", radius=0.1)
    add_text(s, 7.3, 2.55, 1.5, 0.3, "👮 \u95e8\u536b", font_size=11, bold=True, color="9F1239", font=FONT_CN, align="center", valign="middle")
    add_text(s, 7.3, 2.85, 1.5, 0.3, "\u624b\u673a", font_size=9, color="9F1239", font=FONT_CN, align="center", valign="middle")

    # OpenAI external
    add_rect(s, 7.3, 3.4, 1.5, 0.7, fill="EDE9FE", radius=0.1)
    add_text(s, 7.3, 3.4, 1.5, 0.3, "\u2601 OpenAI", font_size=11, bold=True, color="5B21B6", font=FONT_EN, align="center", valign="middle")
    add_text(s, 7.3, 3.7, 1.5, 0.3, "GPT-4o-mini", font_size=9, color="5B21B6", font=FONT_EN, align="center", valign="middle")

    # Deepgram / Edge
    add_rect(s, 7.3, 4.25, 1.5, 0.7, fill="E0F2FE", radius=0.1)
    add_text(s, 7.3, 4.25, 1.5, 0.3, "\u2601 \u5916\u90e8 AI", font_size=11, bold=True, color="1E3A8A", font=FONT_CN, align="center", valign="middle")
    add_text(s, 7.3, 4.55, 1.5, 0.3, "Deepgram + Edge", font_size=9, color="1E3A8A", font=FONT_EN, align="center", valign="middle")

    # Arrows
    add_line(s, 1.8, 2.0, 2.0, 2.0, color="muted", width=1.5)
    add_line(s, 3.3, 2.0, 3.5, 2.0, color="2563EB", width=1.5)
    add_line(s, 7.1, 2.0, 7.3, 2.0, color="16A34A", width=2.0)  # critical
    add_line(s, 7.1, 2.9, 7.3, 2.9, color="muted", width=1.0)

    # Bottom callout: critical design
    add_rect(s, 0.6, 5.0, 6.4, 0.45, fill="FEF3C7", line="F59E0B", radius=0.05)
    add_text(
        s, 0.75, 5.0, 6.2, 0.45,
        "💡 \u5173\u952e\u8bbe\u8ba1\uff1a\u552f\u4e00 `send_to_guard_and_end_call` \u5de5\u5177\u63a7\u5236\u7ed3\u675f\u901a\u8bdd\uff0c\u5fc5\u987b\u5148\u63a8\u9001\u4f01\u4e1a\u5fae\u4fe1\u540e\u624d\u80fd\u6302\u65ad",
        font_size=10, color="92400E", font=FONT_CN, valign="middle",
    )

    add_page_number(s, 6, total)

    # =================== Slide 07 — Section 03 Divider ===================
    s = slide_blank(prs)
    set_bg(s, "bg")
    add_rect(s, 0.6, 0.6, 4.5, 4.4, fill="light", radius=0.1)
    add_text(s, 0.6, 1.1, 4.5, 2.5, "03", font_size=200, bold=True, color="secondary", font=FONT_EN, align="center", valign="middle")
    add_text(s, 0.6, 3.6, 4.5, 0.4, "Section", font_size=14, color="muted", font=FONT_EN, align="center")
    add_text(s, 0.6, 4.0, 4.5, 0.4, "PART THREE", font_size=12, color="muted", font=FONT_EN, align="center", bold=True)

    add_text(s, 5.4, 1.5, 4.2, 0.6, "\u5bf9\u8bdd\u4f53\u9a8c\u6807\u51c6", font_size=36, bold=True, color="primary", font=FONT_CN)
    add_text(s, 5.4, 2.3, 4.2, 0.4, "Conversation Quality Standards", font_size=14, color="secondary", font=FONT_EN)
    add_line(s, 5.4, 2.85, 9.4, 2.85, color="accent", width=2.0)
    add_text(
        s, 5.4, 3.05, 4.2, 1.5,
        "\u4e0d\u8981\u673a\u68b0\u5f0f\u4e00\u95ee\u4e00\u7b54\uff1a\n\u5c11\u8f6e\u6b21 \u00b7 \u591a\u5b57\u6bb5\u5408\u5e76 \u00b7 \u50cf\u771f\u4eba\u95e8\u536b",
        font_size=14, color="muted", font=FONT_CN, line_spacing=1.5,
    )

    add_page_number(s, 7, total)

    # =================== Slide 08 — Conversation comparison ===================
    s = slide_blank(prs)
    set_bg(s, "bg")
    add_text(s, 0.6, 0.4, 9, 0.6, "\u5bf9\u8bdd\u4f53\u9a8c\u5bf9\u6bd4", font_size=28, bold=True, color="primary", font=FONT_CN)
    add_text(s, 0.6, 0.95, 9, 0.4, "Bad vs Good Conversation Examples", font_size=12, color="muted", font=FONT_EN)
    add_line(s, 0.6, 1.35, 9.4, 1.35, color="border", width=1.0)

    # Bad example
    add_rect(s, 0.6, 1.6, 4.4, 3.3, fill="FEF2F2", line="EF4444", radius=0.1)
    add_text(s, 0.8, 1.7, 4.0, 0.4, "\u2717 \u673a\u68b0\u5f0f\u4e00\u95ee\u4e00\u7b54", font_size=14, bold=True, color="EF4444", font=FONT_CN)
    add_text(s, 0.8, 2.1, 4.0, 0.3, "6 \u8f6e \u00b7 \u7ea6 45 s \u00b7 \u5b8c\u5168\u4e0d\u53ef\u63a5\u53d7", font_size=10, color="muted", font=FONT_CN)

    bad_dialog = [
        ("AI: ", "\u8bf7\u95ee\u8f66\u724c\u53f7\uff1f"),
        ("\u7528\u6237: ", "\u6caaA12345"),
        ("AI: ", "\u8bf7\u95ee\u6765\u8bbf\u54ea\u5bb6\u516c\u53f8\uff1f"),
        ("\u7528\u6237: ", "\u84dd\u8272\u9cb8\u9c7c"),
        ("AI: ", "\u8bf7\u95ee\u4e8b\u7531\uff1f"),
        ("\u7528\u6237: ", "\u9001\u8d27"),
        ("AI: ", "\u8bf7\u95ee\u624b\u673a\u53f7\uff1f"),
        ("\u7528\u6237: ", "138xxxx1234"),
        ("AI: ", "\u9884\u8ba1\u505c\u7559\uff1f"),
    ]
    y = 2.5
    for speaker, text in bad_dialog:
        add_text(s, 0.8, y, 4.0, 0.25, [
            (speaker, {"size": 9, "bold": True, "color": "EF4444", "font": FONT_CN}),
            (text, {"size": 9, "color": "primary", "font": FONT_CN}),
        ])
        y += 0.25

    # Good example
    add_rect(s, 5.2, 1.6, 4.4, 3.3, fill="ECFDF5", line="10B981", radius=0.1)
    add_text(s, 5.4, 1.7, 4.0, 0.4, "\u2713 \u81ea\u7136\u5bf9\u8bdd", font_size=14, bold=True, color="10B981", font=FONT_CN)
    add_text(s, 5.4, 2.1, 4.0, 0.3, "3 \u8f6e \u00b7 \u7ea6 15 s \u00b7 \u53c8\u5feb\u53c8\u597d", font_size=10, color="muted", font=FONT_CN)

    good_dialog = [
        ("AI: ", "\u60a8\u597d\uff0c\u8bf7\u8bb2\u4e00\u4e0b\u8f66\u724c\u3001\u6765\u54ea\u5bb6\u516c\u53f8\u3001\u4ec0\u4e48\u4e8b\uff1f"),
        ("\u7528\u6237: ", "\u6caaA12345\uff0c\u6765\u84dd\u8272\u9cb8\u9c7c\u9001\u8d27\u3002"),
        ("AI: ", "\u6536\u5230\uff0c\u624b\u673a\u53f7\u65b9\u4fbf\u7559\u4e00\u4e0b\u5417\uff1f"),
        ("\u7528\u6237: ", "13812345678\u3002"),
        ("AI: ", "\u6caaA12345\uff0c\u84dd\u8272\u9cb8\u9c7c\u9001\u8d27\uff0c\u5df2\u901a\u77e5\u95e8\u536b\uff0c\u8bf7\u7a0d\u7b49\u3002"),
    ]
    y = 2.5
    for speaker, text in good_dialog:
        add_text(s, 5.4, y, 4.0, 0.25, [
            (speaker, {"size": 9, "bold": True, "color": "10B981", "font": FONT_CN}),
            (text, {"size": 9, "color": "primary", "font": FONT_CN}),
        ])
        y += 0.3

    # Bonus: returning
    add_rect(s, 0.6, 5.05, 9.0, 0.4, fill="FEF3C7", line="F59E0B", radius=0.05)
    add_text(
        s, 0.75, 5.05, 8.7, 0.4,
        "\u2728 \u52a0\u5206\uff1a\u56de\u8bbf\u8bc6\u522b\uff0c2 \u8f6e \u00b7 \u7ea6 8 s \u2014\u2014 \u201c\u5f20\u5e08\u5085\u60a8\u597d\uff0c\u4eca\u5929\u662f\u4e0d\u662f\u548c\u4e0a\u6b21\u4e00\u6837\u6765\u9001\u8d27\uff1f\u201d",
        font_size=10, color="92400E", font=FONT_CN, valign="middle",
    )

    add_page_number(s, 8, total)

    # =================== Slide 09 — Latency breakdown ===================
    s = slide_blank(prs)
    set_bg(s, "bg")
    add_text(s, 0.6, 0.4, 9, 0.6, "25 \u79d2 SLA \u65f6\u5ef6\u62c6\u89e3", font_size=28, bold=True, color="primary", font=FONT_CN)
    add_text(s, 0.6, 0.95, 9, 0.4, "End-to-End Latency Breakdown", font_size=12, color="muted", font=FONT_EN)
    add_line(s, 0.6, 1.35, 9.4, 1.35, color="border", width=1.0)

    # Big number
    add_text(s, 0.6, 1.7, 3.0, 0.5, "\u5b9e\u6d4b\u603b\u65f6\u5ef6", font_size=12, color="muted", font=FONT_CN)
    add_text(s, 0.6, 2.1, 3.0, 1.2, "3-6", font_size=80, bold=True, color="10B981", font=FONT_EN)
    add_text(s, 0.6, 3.2, 3.0, 0.4, "\u79d2\uff08\u4e2d\u4f4d\u6570\uff09", font_size=14, color="muted", font=FONT_CN)
    add_text(s, 0.6, 3.7, 3.0, 0.4, "SLA \u4e0a\u9650\uff1a25 \u79d2", font_size=12, color="muted", font=FONT_CN)

    # Stacked horizontal bar
    bar_x = 4.0
    bar_y = 2.0
    bar_w_total = 5.8
    # Each stage: name, latency(s), color, fraction
    stages = [
        ("Twilio \u63a5\u901a + Edge TTS \u5f00\u573a", 0.6, "00B4D8", 0.12),
        ("\u7528\u6237\u8bf4\u8bdd \u00b7 STT \u7ec8\u7a3f", 0.4, "0077B6", 0.08),
        ("LLM \u5355\u8f6e\uff08\u542b function call\uff09", 0.7, "03045E", 0.14),
        ("Edge TTS \u9996\u5305\u00b7\u4e0b\u884c", 0.4, "0077B6", 0.08),
        ("WeChat Webhook \u63a8\u9001", 0.2, "10B981", 0.04),
    ]
    # Total bar = 5s, scale: 5.8" / 5s
    scale = bar_w_total / 5.0
    cur_x = bar_x
    for name, secs, color, _ in stages:
        seg_w = secs * scale
        add_rect(s, cur_x, bar_y, seg_w, 0.6, fill=color, radius=0.03)
        add_text(s, cur_x, bar_y, seg_w, 0.6, f"{secs}s", font_size=10, bold=True, color="white", font=FONT_EN, align="center", valign="middle")
        cur_x += seg_w
    add_text(s, bar_x, bar_y + 0.7, 5.8, 0.3, "\u2192 \u65f6\u95f4\u8f74\uff08\u603b\u8ba1 ~2.3 s\uff09", font_size=10, color="muted", font=FONT_CN, align="center")

    # Why it's fast
    add_text(s, 4.0, 3.0, 5.6, 0.4, "\u4e3a\u4ec0\u4e48\u80fd\u8fbe\u5230 3-6 \u79d2\uff1f", font_size=14, bold=True, color="primary", font=FONT_CN)
    reasons = [
        "1. \u6d41\u5f0f STT\uff1a\u7528\u6237\u8bf4\u5b8c\u524d\u5df2\u7ecf\u62ff\u5230 interim \u7ed3\u679c",
        "2. \u589e\u91cf LLM\uff1a\u542c\u5230\u4e00\u4e2a\u5b57\u6bb5\u5c31\u8c03 update_visitor_info",
        "3. \u6d41\u5f0f TTS\uff1a\u7b2c\u4e00\u4e2a\u53e5\u5b50\u751f\u6210\u540e\u7acb\u5373\u4e0b\u884c\uff0c\u540e\u7eed\u540c\u6b65\u751f\u6210",
        "4. \u4e0a\u4e0b\u6587\u538b\u7f29\uff1a\u7cfb\u7edf\u63d0\u793a\u8bcd < 1.1 KB\uff0cTTFT < 600 ms",
        "5. \u51fd\u6570\u8c03\u7528\u5355\u70b9\u6536\u53e3\uff1aWeChat \u63a8\u9001\u4e0e LLM \u5e94\u7b54\u5e76\u884c",
    ]
    y = 3.5
    for r in reasons:
        add_text(s, 4.0, y, 5.6, 0.3, r, font_size=11, color="primary", font=FONT_CN)
        y += 0.3

    add_page_number(s, 9, total)

    # =================== Slide 10 — WeChat integration + closing ===================
    s = slide_blank(prs)
    set_bg(s, "bg")
    add_text(s, 0.6, 0.4, 9, 0.6, "\u4f01\u4e1a\u5fae\u4fe1\u96c6\u6210\u4e0e\u4ea4\u4ed8\u6e05\u5355", font_size=28, bold=True, color="primary", font=FONT_CN)
    add_text(s, 0.6, 0.95, 9, 0.4, "WeChat Work Integration & Delivery Checklist", font_size=12, color="muted", font=FONT_EN)
    add_line(s, 0.6, 1.35, 9.4, 1.35, color="border", width=1.0)

    # Left: WeChat flow
    add_text(s, 0.6, 1.55, 4.4, 0.4, "💬 \u4f01\u4e1a\u5fae\u4fe1\u96c6\u6210", font_size=16, bold=True, color="primary", font=FONT_CN)
    wechat_steps = [
        ("\u7fa4\u673a\u5668\u4eba Webhook", "5 \u5206\u949f\u63a5\u5165\uff0c0 \u5ba1\u6279\uff0c0 \u8d39\u7528"),
        ("Markdown \u5361\u7247\u63a8\u9001", "\u8f66\u724c\u00b7\u516c\u53f8\u00b7\u4e8b\u7531\u00b7\u624b\u673a\u00b7\u9884\u8ba1\u505c\u7559"),
        ("\u5386\u53f2\u6765\u8bbf\u5217\u8868", "\u56de\u8bbf\u8bc6\u522b\u4e0e\u53ef\u5ba1\u8ba1"),
        ("\u9519\u8bef\u91cd\u8bd5", "errcode \u4f18\u5148\u7ea7\u5224\u65ad\uff0c\u4e0d\u91cd\u8bd5\u6388\u6743\u9519"),
        ("\u4ea4\u4ed8\u5f62\u5f0f", "\u4ee3\u7801\u5b8c\u6574 + .env.example + README + \u67b6\u6784\u56fe"),
    ]
    y = 2.0
    for name, desc in wechat_steps:
        add_rect(s, 0.6, y, 0.08, 0.4, fill="00B4D8")
        add_text(s, 0.8, y, 1.6, 0.4, name, font_size=11, bold=True, color="primary", font=FONT_CN, valign="middle")
        add_text(s, 2.4, y, 2.6, 0.4, desc, font_size=10, color="muted", font=FONT_CN, valign="middle")
        y += 0.45

    # Right: WeChat preview
    add_text(s, 5.4, 1.55, 4.2, 0.4, "📸 \u4f01\u4e1a\u5fae\u4fe1\u5361\u7247\u9884\u89c8", font_size=16, bold=True, color="primary", font=FONT_CN)
    add_rect(s, 5.4, 2.0, 4.2, 3.05, fill="white", line="E2E8F0", radius=0.1)
    add_text(s, 5.55, 2.1, 4.0, 0.4, "## 🚗 \u65b0\u8bbf\u5ba2\u767b\u8bb0 🆕 \u56de\u8bbf", font_size=11, bold=True, color="primary", font=FONT_CN)
    preview = [
        ("\u8f66\u724c\u53f7\uff1a", "\u6caaA12345"),
        ("\u53d7\u8bbf\u516c\u53f8\uff1a", "\u84dd\u8272\u9cb8\u9c7c\u79d1\u6280"),
        ("\u6765\u8bbf\u4e8b\u7531\uff1a", "\u9001\u8d27"),
        ("\u8bbf\u5ba2\u59d3\u540d\uff1a", "\u5f20\u5e08\u5085"),
        ("\u8054\u7cfb\u7535\u8bdd\uff1a", "13812345678"),
        ("\u9884\u8ba1\u505c\u7559\uff1a", "2 \u5c0f\u65f6"),
    ]
    y = 2.55
    for k, v in preview:
        add_text(s, 5.55, y, 1.5, 0.3, k, font_size=10, color="muted", font=FONT_CN)
        add_text(s, 7.05, y, 2.4, 0.3, v, font_size=10, bold=True, color="primary", font=FONT_CN)
        y += 0.3
    add_text(s, 5.55, 4.4, 4.0, 0.3, "\u2705 \u95e8\u536b\u8bf7\u7acb\u5373\u653e\u884c", font_size=11, bold=True, color="10B981", font=FONT_CN)
    add_text(s, 5.55, 4.7, 4.0, 0.3, "_ \u901a\u8bdd\u65f6\u957f 12.4s _", font_size=9, color="muted", font=FONT_CN)

    # Bottom: thanks
    add_rect(s, 0.6, 5.1, 9.0, 0.35, fill="light", radius=0.05)
    add_text(
        s, 0.6, 5.1, 9.0, 0.35,
        "\u8c22\u8c22 \u00b7 Thanks for your time \u00b7 hr@whaletech.ai",
        font_size=11, bold=True, color="secondary", font=FONT_CN, align="center", valign="middle",
    )

    add_page_number(s, 10, total)

    # =================== Save ===================
    out_dir = Path(__file__).resolve().parent.parent / "docs" / "pptx" / "output"
    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / "voice-agent-take-home.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")
    return out


if __name__ == "__main__":
    build()
